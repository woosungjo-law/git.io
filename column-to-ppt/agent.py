import os
import sys
import json
import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


def read_txt_file(filepath: str) -> str:
    with open(filepath, "r", encoding="utf-8") as f:
        return f.read()


def ask_slide_count(client: anthropic.Anthropic, text: str) -> int:
    print("\n[AI]: 텍스트를 읽었습니다. 몇 장의 슬라이드로 만들어 드릴까요?")
    while True:
        user_input = input("[사용자]: ").strip()
        if user_input.isdigit() and int(user_input) > 0:
            return int(user_input)
        print("[AI]: 숫자로 입력해 주세요. (예: 5)")


def generate_slides_content(client: anthropic.Anthropic, text: str, slide_count: int) -> list[dict]:
    print(f"\n[AI]: {slide_count}장 슬라이드 내용을 생성 중입니다...")

    prompt = f"""다음 텍스트를 분석하여 {slide_count}장의 PPT 슬라이드 내용을 만들어 주세요.

텍스트:
{text}

다음 JSON 형식으로 정확히 {slide_count}개의 슬라이드를 반환하세요:
{{
  "slides": [
    {{
      "title": "슬라이드 제목",
      "bullets": ["핵심 내용 1", "핵심 내용 2", "핵심 내용 3"]
    }}
  ]
}}

규칙:
- 첫 번째 슬라이드는 전체 내용의 제목/소개
- 마지막 슬라이드는 요약 또는 결론
- 각 슬라이드의 bullets는 2~4개
- 한국어로 작성
- JSON만 반환, 다른 텍스트 없음"""

    response = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=4096,
        messages=[{"role": "user", "content": prompt}],
    )

    raw = response.content[0].text.strip()
    # JSON 블록 추출
    if "```" in raw:
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]

    data = json.loads(raw)
    return data["slides"]


def create_pptx(slides: list[dict], output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for i, slide_data in enumerate(slides):
        if i == 0:
            layout = prs.slide_layouts[0]  # 제목 슬라이드
        else:
            layout = prs.slide_layouts[1]  # 제목 + 내용

        slide = prs.slides.add_slide(layout)

        # 제목
        title = slide.shapes.title
        title.text = slide_data["title"]
        title.text_frame.paragraphs[0].font.size = Pt(32)
        title.text_frame.paragraphs[0].font.bold = True
        title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x1F, 0x49, 0x7D)

        # 내용 (bullets)
        if i > 0 and len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for j, bullet in enumerate(slide_data.get("bullets", [])):
                if j == 0:
                    tf.paragraphs[0].text = bullet
                    tf.paragraphs[0].font.size = Pt(20)
                else:
                    p = tf.add_paragraph()
                    p.text = bullet
                    p.font.size = Pt(20)

    prs.save(output_path)
    print(f"\n[AI]: PPT 파일이 생성되었습니다 → {output_path}")


def main():
    if len(sys.argv) < 2:
        print("사용법: python3 agent.py <텍스트파일.txt>")
        sys.exit(1)

    txt_path = sys.argv[1]
    if not os.path.exists(txt_path):
        print(f"파일을 찾을 수 없습니다: {txt_path}")
        sys.exit(1)

    api_key = os.environ.get("ANTHROPIC_API_KEY")
    if not api_key:
        print("ANTHROPIC_API_KEY 환경변수를 설정해 주세요.")
        sys.exit(1)

    client = anthropic.Anthropic(api_key=api_key)

    text = read_txt_file(txt_path)
    print(f"[AI]: '{txt_path}' 파일을 읽었습니다. ({len(text)}자)")

    slide_count = ask_slide_count(client, text)
    slides = generate_slides_content(client, text, slide_count)

    output_path = txt_path.replace(".txt", ".pptx")
    if output_path == txt_path:
        output_path = txt_path + ".pptx"

    create_pptx(slides, output_path)


if __name__ == "__main__":
    main()
